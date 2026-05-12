import os
import re
from pathlib import Path
from typing import Optional

from docx import Document


class FileCreator:
    def __init__(self) -> None:
        self.generated = Path("generated_files")
        self.generated.mkdir(exist_ok=True)

    def create_txt(self, name: str, content: str, folder: Optional[str] = None) -> str:
        path = self._target_path(name, ".txt", folder)
        path.write_text(content, encoding="utf-8")
        return str(path.resolve())

    def create_docx(self, name: str, content: str, folder: Optional[str] = None) -> str:
        path = self._target_path(name, ".docx", folder)
        doc = Document()
        title = path.stem.replace("_", " ").strip()
        if title:
            doc.add_heading(title, 0)
        for block in content.split("\n"):
            clean = block.strip()
            if not clean:
                continue
            if clean.startswith("# "):
                doc.add_heading(clean[2:], level=1)
            elif clean.startswith("## "):
                doc.add_heading(clean[3:], level=2)
            else:
                doc.add_paragraph(clean)
        doc.save(path)
        return str(path.resolve())

    def create_pptx(self, name: str, content: str, folder: Optional[str] = None) -> str:
        from pptx import Presentation
        path = self._target_path(name, ".pptx", folder)
        prs = Presentation()
        sections = [s.strip() for s in re.split(r"\n\s*---\s*\n", content) if s.strip()]
        if not sections:
            sections = [content]
        for section in sections:
            lines = [l.strip() for l in section.splitlines() if l.strip()]
            title = lines[0].lstrip("# ") if lines else "Slide"
            body = "\n".join(lines[1:]) if len(lines) > 1 else ""
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            slide.placeholders[1].text = body
        prs.save(path)
        return str(path.resolve())

    def _target_path(self, name: str, ext: str, folder: Optional[str]) -> Path:
        safe = self._safe(name)
        if not safe.lower().endswith(ext):
            safe += ext
        if folder:
            base = self._folder(folder)
        else:
            base = self.generated
        base.mkdir(parents=True, exist_ok=True)
        path = base / safe
        i = 2
        while path.exists():
            path = base / f"{Path(safe).stem}_{i}{ext}"
            i += 1
        return path

    def _safe(self, name: str) -> str:
        name = name.strip().replace(" ", "_")
        name = re.sub(r"[^a-zA-Z0-9_\-\.åäöÅÄÖ]+", "", name)
        return name or "jarvis_file"

    def _folder(self, name: str) -> Path:
        key = name.lower().strip()
        home = Path.home()
        mapping = {
            "desktop": home / "Desktop",
            "downloads": home / "Downloads",
            "documents": home / "Documents",
            "generated": self.generated,
        }
        return mapping.get(key, self.generated)
